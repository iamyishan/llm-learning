import json
from langchain_community.llms.ollama import Ollama
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from pptx import Presentation


# 调用阿里云百炼大模型API
def ollama_llm(query, history=[], user_stop_words=[]):
    llm = Ollama(base_url="http://localhost:8000", model="qwen2:7b", verbose=True)

    prompt = ChatPromptTemplate.from_messages([
        ("system", "你是一个PPT制作小助手"),
        ("user", "{input}")
    ])
    output_parser = StrOutputParser()
    chain = prompt | llm | output_parser
    return chain.invoke({"input": query})


# 生成PPT内容
def generate_ppt_content(topic, pages):
    # 输出格式
    output_format = json.dumps({
        "title": "example title",
        "pages": [
            {
                "title": "title for page 1",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                ],
            },
            {
                "title": "title for page 2",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                    {
                        "title": "title for paragraph 3",
                        "description": "detail for paragraph 3",
                    },
                ],
            },
        ],
    }, ensure_ascii=True)

    # prompt
    prompt = f'''我要准备1个关于{topic}的PPT，要求一共写{pages}页，请你根据主题生成详细内容，不要省略。
    按这个JSON格式输出{output_format}，只能返回JSON，且JSON不要用```包裹，内容要用中文。'''

    print(prompt)
    print("=============================")
    resp = ollama_llm(prompt)
    print(resp)
    # 调用llm生成PPT内容
    ppt_content = json.loads(resp)
    return ppt_content


# 生成PPT文件
def generate_ppt_file(topic, ppt_content):
    ppt = Presentation()

    # PPT首页
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # title&subtitle layout
    slide.placeholders[0].text = ppt_content['title']
    slide.placeholders[1].text = "通义千问72B"

    # 内容页
    print('总共%d页...' % len(ppt_content['pages']))
    for i, page in enumerate(ppt_content['pages']):
        print('生成第%d页:%s' % (i + 1, page['title']))
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # title&content layout
        # 标题
        slide.placeholders[0].text = page['title']
        # 正文
        for sub_content in page['content']:
            print(sub_content)
            # 一级正文
            sub_title = slide.placeholders[1].text_frame.add_paragraph()
            sub_title.text, sub_title.level = sub_content['title'], 1
            # 二级正文
            sub_description = slide.placeholders[1].text_frame.add_paragraph()
            sub_description.text, sub_description.level = sub_content['description'], 2

    ppt.save('%s.pptx' % topic)


if __name__ == '__main__':
    # query = "请帮我写一个关于三维重建的PPT，一共写5页"
    # ollama_llm(query)
    while True:
        # 输入需求
        topic = input('输入主题:')
        pages = int(input('输入页数:'))
        # 生成PPT内容
        ppt_content = generate_ppt_content(topic, pages)
        # 生成PPT文件
        generate_ppt_file(topic, ppt_content)
